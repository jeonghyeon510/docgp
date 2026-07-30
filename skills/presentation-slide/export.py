#!/usr/bin/env python3
"""HTML 슬라이드를 PPTX / PDF로 내보낸다.

동작 방식
    1. 설치된 Chrome(또는 Edge/Chromium)을 헤드리스로 띄워 슬라이드를 한 장씩 PNG로 캡처
    2. python-pptx로 16:9 프레젠테이션을 만들고 PNG를 한 장씩 꽉 채워 배치
    3. HTML 안의 <aside class="notes">를 읽어 각 슬라이드의 '발표자 노트'로 삽입
    4. 같은 PNG들을 Pillow로 묶어 PDF 생성

필요한 것
    · Chrome / Edge / Chromium 중 하나 (대부분 이미 설치되어 있음)
    · python-pptx, Pillow      →  python3 -m pip install python-pptx Pillow

사용법
    python3 export.py slide.html                  # PPTX + PDF 둘 다
    python3 export.py slide.html --only pptx      # PPTX만
    python3 export.py slide.html --only pdf       # PDF만
    python3 export.py slide.html --scale 1        # 캡처 배율(기본 2 = 고화질)
    python3 export.py slide.html -o out/발표자료   # 출력 파일명 지정(확장자 제외)
"""
from __future__ import annotations

import argparse
import glob
import html
import os
import platform
import re
import shutil
import subprocess
import sys
from pathlib import Path
from urllib.parse import unquote

SLIDE_W, SLIDE_H = 1280, 720  # slide.html의 --w/--h와 반드시 같아야 한다


# ─────────────────────────────────────────────────────────────
# 1. 브라우저 찾기
# ─────────────────────────────────────────────────────────────
def find_browser() -> str:
    """Chrome/Edge/Chromium 실행 파일 경로를 찾는다.

    playwright가 받아둔 chrome-headless-shell이 있으면 그걸 우선 쓴다.
    캡처 전용으로 만들어진 빌드라 정식 Chrome보다 가볍고 확실하게 종료된다.
    """
    system = platform.system()
    candidates: list[str] = []

    def cache(*patterns: str) -> list[str]:
        found: list[str] = []
        for pattern in patterns:
            found += sorted(glob.glob(os.path.expanduser(os.path.expandvars(pattern))), reverse=True)
        return found

    if system == "Darwin":
        candidates += cache("~/Library/Caches/ms-playwright/chromium_headless_shell-*"
                            "/chrome-headless-shell-mac-*/chrome-headless-shell")
        candidates += [
            "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
            "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
            "/Applications/Chromium.app/Contents/MacOS/Chromium",
        ]
        candidates += cache("~/Library/Caches/ms-playwright/chromium-*/chrome-mac/Chromium.app"
                            "/Contents/MacOS/Chromium")
    elif system == "Windows":
        candidates += cache(r"%LOCALAPPDATA%\ms-playwright\chromium_headless_shell-*"
                            r"\chrome-headless-shell-win*\chrome-headless-shell.exe")
        for root in (os.environ.get("PROGRAMFILES", r"C:\Program Files"),
                     os.environ.get("PROGRAMFILES(X86)", r"C:\Program Files (x86)"),
                     os.environ.get("LOCALAPPDATA", "")):
            if not root:
                continue
            candidates += [
                rf"{root}\Google\Chrome\Application\chrome.exe",
                rf"{root}\Microsoft\Edge\Application\msedge.exe",
            ]
        candidates += cache(r"%LOCALAPPDATA%\ms-playwright\chromium-*\chrome-win\chrome.exe")
    else:  # Linux
        candidates += cache("~/.cache/ms-playwright/chromium_headless_shell-*"
                            "/chrome-headless-shell-linux*/chrome-headless-shell")
        for name in ("google-chrome", "chromium", "chromium-browser", "microsoft-edge"):
            found = shutil.which(name)
            if found:
                candidates.append(found)
        candidates += cache("~/.cache/ms-playwright/chromium-*/chrome-linux/chrome")

    for path in candidates:
        if path and os.path.exists(path):
            return path

    sys.exit(
        "Chrome/Edge/Chromium을 찾지 못했습니다.\n"
        "  · Windows: Edge가 기본 설치되어 있어야 정상입니다\n"
        "  · macOS  : https://google.com/chrome 에서 Chrome을 설치하세요\n"
        "  · 이미 설치돼 있다면 CHROME_PATH 환경변수로 경로를 직접 지정하세요"
    )


# ─────────────────────────────────────────────────────────────
# 2. HTML에서 발표자 노트 뽑기
# ─────────────────────────────────────────────────────────────
SECTION_RE = re.compile(r'<section[^>]*\bclass="[^"]*\bslide\b[^"]*"[^>]*>(.*?)</section>',
                        re.S | re.I)
NOTES_RE = re.compile(r'<aside[^>]*\bclass="[^"]*\bnotes\b[^"]*"[^>]*>(.*?)</aside>', re.S | re.I)
TAG_RE = re.compile(r"<[^>]+>")
# 주석 안에 설명용으로 적어둔 <section class="slide"> 같은 문자열을 실제 태그로 오인하지 않도록
# 파싱 전에 HTML 주석을 통째로 걷어낸다. (발표자 노트는 주석이 아니라 <aside>라 영향 없음)
COMMENT_RE = re.compile(r"<!--.*?-->", re.S)


def read_sections(html_path: Path) -> list[str]:
    source = COMMENT_RE.sub("", html_path.read_text(encoding="utf-8"))
    return SECTION_RE.findall(source)


def extract_notes(html_path: Path) -> list[str]:
    """슬라이드 순서대로 발표자 노트 텍스트를 반환한다 (없으면 빈 문자열)."""
    notes: list[str] = []
    for section in read_sections(html_path):
        match = NOTES_RE.search(section)
        if not match:
            notes.append("")
            continue
        text = TAG_RE.sub("", match.group(1))
        notes.append(html.unescape(text).strip())
    return notes


def count_slides(html_path: Path) -> int:
    return len(read_sections(html_path))


# ─────────────────────────────────────────────────────────────
# 2-1. 이미지 경로 점검
# ─────────────────────────────────────────────────────────────
IMG_RE = re.compile(r"""<img[^>]*\bsrc\s*=\s*["']([^"']+)["']""", re.I)
PLACEHOLDER_RE = re.compile(r"""class\s*=\s*["'][^"']*\bplaceholder\b""", re.I)
WINDOWS_DRIVE_RE = re.compile(r"^[A-Za-z]:")


def check_images(html_path: Path) -> "tuple[list[str], list[str]]":
    """<img>가 가리키는 파일이 실제로 있는지, 다른 컴퓨터에서도 열릴 경로인지 본다.

    이미지가 없거나 경로가 이 컴퓨터에서만 통하는 형태면 브라우저는 아무 경고 없이
    빈 칸으로 그리고, 그 빈 칸이 그대로 PPTX에 박힌다. 그래서 캡처 전에 먼저 막는다.

    반환값은 (중단할 문제, 경고)다.
    """
    base = html_path.resolve().parent
    source = COMMENT_RE.sub("", html_path.read_text(encoding="utf-8"))
    errors: list[str] = []

    for src in IMG_RE.findall(source):
        # 인라인(data:)과 외부(http:) 이미지는 로컬 파일 검사 대상이 아니다
        if src.startswith("data:") or re.match(r"https?://", src, re.I):
            continue

        if "\\" in src:
            errors.append(f"{src}\n      → 역슬래시(\\)가 섞여 있다. 윈도우에서 작업했더라도 HTML 경로는 항상 슬래시(/)를 쓴다")
            continue
        if src.startswith("file://") or src.startswith("/") or WINDOWS_DRIVE_RE.match(src):
            errors.append(f"{src}\n      → 절대경로라 다른 컴퓨터에서 깨진다. images/파일명 처럼 상대경로로 바꾼다")
            continue

        # macOS·윈도우는 파일명 대소문자를 가리지 않아 exists()가 그냥 통과시킨다.
        # 그래서 존재 여부만 믿지 않고 디스크에 적힌 실제 이름과 철자까지 맞춰본다.
        # (대소문자가 어긋난 채로 넘어가면 다른 환경에서 그림만 안 나온다)
        target = base / unquote(src)
        actual = None
        if target.parent.is_dir():
            actual = next((p for p in target.parent.iterdir()
                           if p.name.lower() == target.name.lower()), None)

        if actual is None:
            errors.append(f"{src}\n      → 파일이 없다. 여기에 있어야 한다: {target}")
        elif actual.name != target.name:
            errors.append(f"{src}\n      → 대소문자가 다르다. 실제 파일명은 '{actual.name}'")

    warnings = []
    if PLACEHOLDER_RE.search(source):
        warnings.append("점선 자리표시(.placeholder) 박스가 아직 남아 있다 — 이미지를 넣고 지울 것")

    return errors, warnings


# ─────────────────────────────────────────────────────────────
# 3. 슬라이드 캡처
# ─────────────────────────────────────────────────────────────
def capture(browser: str, html_path: Path, total: int, out_dir: Path, scale: int) -> list[Path]:
    """슬라이드를 한 장씩 PNG로 캡처해 경로 목록을 돌려준다."""
    url_base = html_path.resolve().as_uri()
    images: list[Path] = []

    for index in range(1, total + 1):
        png = out_dir / f"slide-{index:03d}.png"
        # ⚠ --virtual-time-budget 을 넣으면 Chrome이 종료되지 않고 멈춘다. 절대 추가하지 말 것.
        cmd = [
            browser,
            "--headless",
            "--disable-gpu",
            "--hide-scrollbars",
            "--no-sandbox",
            "--allow-file-access-from-files",
            f"--force-device-scale-factor={scale}",
            f"--window-size={SLIDE_W},{SLIDE_H}",
            f"--screenshot={png}",
            f"{url_base}?export=1&slide={index}",
        ]
        try:
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=90)
        except subprocess.TimeoutExpired:
            sys.exit(f"슬라이드 {index} 캡처가 90초를 넘겨 중단했습니다. 브라우저: {browser}")

        if not png.exists():
            sys.exit(f"슬라이드 {index} 캡처 실패\n{result.stderr[-1500:]}")
        images.append(png)
        print(f"  캡처 {index}/{total}", end="\r", flush=True)

    print(f"  캡처 완료 ({total}장)          ")
    return images


# ─────────────────────────────────────────────────────────────
# 4. PPTX / PDF 만들기
# ─────────────────────────────────────────────────────────────
def build_pptx(images: list[Path], notes: list[str], out: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]       # 완전히 빈 레이아웃

    for index, image in enumerate(images):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
        note = notes[index] if index < len(notes) else ""
        if note:
            slide.notes_slide.notes_text_frame.text = note

    prs.save(out)
    filled = sum(1 for n in notes if n)
    print(f"  PPTX 생성: {out}  (발표자 노트 {filled}장)")


def build_pdf(images: list[Path], out: Path, scale: int) -> None:
    from PIL import Image

    # 해상도를 맞춰야 PDF 페이지가 13.333 x 7.5 인치로 떨어진다
    dpi = SLIDE_W * scale / 13.333
    frames = [Image.open(p).convert("RGB") for p in images]
    frames[0].save(out, save_all=True, append_images=frames[1:], resolution=dpi)
    print(f"  PDF 생성:  {out}  ({len(frames)}페이지)")


# ─────────────────────────────────────────────────────────────
def main() -> None:
    parser = argparse.ArgumentParser(description="HTML 슬라이드를 PPTX/PDF로 내보낸다")
    parser.add_argument("html", help="슬라이드 HTML 파일")
    parser.add_argument("-o", "--out", help="출력 파일명(확장자 제외). 기본값은 HTML과 같은 이름")
    parser.add_argument("--only", choices=["pptx", "pdf", "png"], help="하나만 생성")
    parser.add_argument("--scale", type=int, default=2, help="캡처 배율 (기본 2)")
    parser.add_argument("--keep-png", action="store_true", help="중간 PNG를 지우지 않는다")
    parser.add_argument("--skip-image-check", action="store_true",
                        help="이미지 경로 점검을 건너뛴다 (문제를 알면서 진행할 때만)")
    args = parser.parse_args()

    html_path = Path(args.html)
    if not html_path.exists():
        sys.exit(f"파일을 찾을 수 없습니다: {html_path}")

    total = count_slides(html_path)
    if total == 0:
        sys.exit('슬라이드를 찾지 못했습니다. <section class="slide"> 형식인지 확인하세요.')

    errors, warnings = check_images(html_path)
    for warning in warnings:
        print(f"경고: {warning}")
    if errors and not args.skip_image_check:
        print(f"\n이미지 {len(errors)}개에 문제가 있습니다. 이대로 내보내면 빈 칸으로 박힙니다.\n")
        for problem in errors:
            print(f"  · {problem}")
        sys.exit("\n고친 뒤 다시 실행하세요. 확인했고 그냥 진행하려면 --skip-image-check 를 붙이세요.")

    out_base = Path(args.out) if args.out else html_path.with_suffix("")
    out_base.parent.mkdir(parents=True, exist_ok=True)

    browser = os.environ.get("CHROME_PATH") or find_browser()
    print(f"브라우저: {browser}")
    print(f"슬라이드: {total}장")

    png_dir = out_base.parent / f"{out_base.name}-png"
    png_dir.mkdir(exist_ok=True)

    images = capture(browser, html_path, total, png_dir, args.scale)

    if args.only != "png":
        notes = extract_notes(html_path)
        if args.only in (None, "pptx"):
            build_pptx(images, notes, out_base.with_suffix(".pptx"))
        if args.only in (None, "pdf"):
            build_pdf(images, out_base.with_suffix(".pdf"), args.scale)

    if args.only == "png" or args.keep_png:
        print(f"  PNG 유지:  {png_dir}/")
    else:
        shutil.rmtree(png_dir, ignore_errors=True)

    print("완료.")


if __name__ == "__main__":
    main()
